"""Local PDF -> editable PPTX converter.

Ported from ysrock/pdf2pptx-ai-tool with four targeted fixes:
  1. Slide canvas — fixed 16:9, each page image fit inside (letterbox)
  2. Default DPI 200 (was 100) for better OCR
  3. OCR noise filtering (low score, single chars, tiny bboxes)
  4. Font size = bbox_h * 0.65 (was 0.8 — too large)
"""
from __future__ import annotations

import gc
import json
import threading
import time
from io import BytesIO
from pathlib import Path
from typing import Callable, Optional

import cv2
import numpy as np
import pdfplumber
import torch
from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt
from rapidocr_onnxruntime import RapidOCR
from simple_lama_inpainting import SimpleLama


# Fix #1: every slide is 16:9 widescreen, EMU = 914400 per inch.
SLIDE_W_EMU = Emu(int(13.333 * 914400))
SLIDE_H_EMU = Emu(int(7.5 * 914400))
SLIDE_W_PT = 13.333 * 72
SLIDE_H_PT = 7.5 * 72

# Fix #3 thresholds
MIN_OCR_SCORE = 0.5
MIN_TEXT_LEN = 2  # at least 2 chars unless korean (single hangul ok)
MIN_BBOX_AREA_RATIO = 1e-4  # bbox_w * bbox_h / page_w * page_h

# pdfplumber uses pypdfium2 under the hood, which is NOT thread-safe.
# We serialize all pdfplumber.open() calls in this module so that the
# conversion thread doesn't race with itself. The thumbnail path uses
# PyMuPDF instead and is independent.
_pdf_lock = threading.Lock()


def _is_korean(s: str) -> bool:
    return any("가" <= ch <= "힣" for ch in s)


def _bbox_to_rect(box) -> tuple[float, float, float, float]:
    xs = [p[0] for p in box]
    ys = [p[1] for p in box]
    return min(xs), min(ys), max(xs), max(ys)


def _bbox_iou(box1, box2) -> float:
    """Axis-aligned IoU on quadrilateral bboxes (uses bounding rectangle)."""
    x1a, y1a, x2a, y2a = _bbox_to_rect(box1)
    x1b, y1b, x2b, y2b = _bbox_to_rect(box2)
    ix1 = max(x1a, x1b)
    iy1 = max(y1a, y1b)
    ix2 = min(x2a, x2b)
    iy2 = min(y2a, y2b)
    if ix2 < ix1 or iy2 < iy1:
        return 0.0
    inter = (ix2 - ix1) * (iy2 - iy1)
    a = max((x2a - x1a) * (y2a - y1a), 1.0)
    b = max((x2b - x1b) * (y2b - y1b), 1.0)
    return inter / (a + b - inter)


def _dedupe_against(new_results: list, existing_results: list, iou_thresh: float = 0.3) -> list:
    """Drop new_results that overlap (>iou_thresh) with any existing one."""
    kept = []
    for r in new_results:
        if any(_bbox_iou(r[0], e[0]) > iou_thresh for e in existing_results):
            continue
        kept.append(r)
    return kept


def _filter_ocr_results(ocr_result, page_w_px: int, page_h_px: int) -> list:
    """Drop noisy OCR results (Fix #3)."""
    if not ocr_result:
        return []
    page_area = max(page_w_px * page_h_px, 1)
    kept = []
    for item in ocr_result:
        try:
            box, text, score = item
        except (ValueError, TypeError):
            continue
        if not text or not text.strip():
            continue
        if score < MIN_OCR_SCORE:
            continue
        text_clean = text.strip()
        # Allow short text only if Korean or numeric/important
        if len(text_clean) < MIN_TEXT_LEN and not _is_korean(text_clean):
            continue
        xs = [p[0] for p in box]
        ys = [p[1] for p in box]
        bw = max(xs) - min(xs)
        bh = max(ys) - min(ys)
        if bw <= 0 or bh <= 0:
            continue
        if (bw * bh) / page_area < MIN_BBOX_AREA_RATIO:
            continue
        kept.append((box, text_clean, score))
    return kept


class LocalConverter:
    """OCR + inpainting pipeline. Initialize models once, reuse across pages."""

    def __init__(self) -> None:
        self.ocr_engine: Optional[RapidOCR] = None
        self.lama: Optional[SimpleLama] = None
        self.is_gpu = False
        self._initialized = False
        self._init_lock = threading.Lock()

    def initialize_models(self) -> None:
        with self._init_lock:
            if self._initialized:
                return
            self.ocr_engine = RapidOCR()
            self.is_gpu = torch.cuda.is_available()
            if not self.is_gpu:
                # Force CPU map_location for SimpleLama's torch.jit.load
                original_jit_load = torch.jit.load
                def safe_jit_load(*args, **kwargs):
                    if "map_location" not in kwargs:
                        kwargs["map_location"] = torch.device("cpu")
                    return original_jit_load(*args, **kwargs)
                torch.jit.load = safe_jit_load
                try:
                    self.lama = SimpleLama()
                finally:
                    torch.jit.load = original_jit_load
            else:
                self.lama = SimpleLama()
            self._initialized = True

    # ------- Per-page processing -------

    def process_page(
        self,
        pil_image: Image.Image,
        dpi: int,
        dilation_size: int,
    ) -> tuple[Image.Image, list[dict]]:
        """Run OCR + inpainting on a PIL image. Returns (clean_bg, text_blocks)."""
        if not self._initialized:
            self.initialize_models()
        assert self.ocr_engine is not None and self.lama is not None

        img_np = np.array(pil_image)
        page_h_px, page_w_px = img_np.shape[:2]

        # ---- Pass 1: standard OCR ----
        ocr_result, _ = self.ocr_engine(img_np)
        ocr_filtered = _filter_ocr_results(ocr_result, page_w_px, page_h_px)

        # ---- Fix B: Pass 2 — whiten the regions Pass 1 covered, then OCR
        # the remainder. RapidOCR re-detection on a sparser image often
        # picks up small/stylised text that Pass 1 missed (chart labels,
        # arrow callouts, decorative headings).
        if ocr_filtered:
            cover = np.zeros(img_np.shape[:2], dtype=np.uint8)
            for box, _t, _s in ocr_filtered:
                poly = np.array(box, dtype=np.float32)
                ctr = np.mean(poly, axis=0)
                expanded = ctr + (poly - ctr) * 1.20  # 20% expansion to catch tails
                pts = expanded.astype(np.int32).reshape((-1, 1, 2))
                cv2.fillPoly(cover, [pts], 255)
            img_pass2 = img_np.copy()
            img_pass2[cover > 0] = 255  # whiten covered regions
            try:
                ocr_result_2, _ = self.ocr_engine(img_pass2)
                ocr_filtered_2 = _filter_ocr_results(ocr_result_2, page_w_px, page_h_px)
                # Drop pass-2 hits that overlap an existing pass-1 hit.
                ocr_filtered_2 = _dedupe_against(ocr_filtered_2, ocr_filtered, iou_thresh=0.2)
                ocr_filtered = ocr_filtered + ocr_filtered_2
            except Exception:
                pass  # Pass 2 is best-effort; never fail the page over it.

        mask = np.zeros(img_np.shape[:2], dtype=np.uint8)
        text_blocks: list[dict] = []
        box_heights: list[float] = []
        scale_factor = 72.0 / dpi  # px -> pt

        for box, text, score in ocr_filtered:
            poly = np.array(box, dtype=np.float32)
            center = np.mean(poly, axis=0)
            expanded_poly = center + (poly - center) * 1.10
            pts = expanded_poly.astype(np.int32).reshape((-1, 1, 2))
            cv2.fillPoly(mask, [pts], 255)

            xs = [p[0] for p in box]
            ys = [p[1] for p in box]
            text_blocks.append({
                "text": text,
                "x_pt": min(xs) * scale_factor,
                "y_pt": min(ys) * scale_factor,
                "w_pt": (max(xs) - min(xs)) * scale_factor,
                "h_pt": (max(ys) - min(ys)) * scale_factor,
                "score": score,
            })
            box_heights.append(max(ys) - min(ys))

        # Adaptive dilation
        if box_heights:
            avg_h = sum(box_heights) / len(box_heights)
            adaptive_dil = int(avg_h * 0.20)
            adaptive_dil = max(15, min(adaptive_dil, 40))
            final_dil = max(dilation_size, adaptive_dil)
            kernel_close = np.ones((5, 5), np.uint8)
            mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, kernel_close)
            kernel = np.ones((final_dil, final_dil), np.uint8)
            mask = cv2.dilate(mask, kernel, iterations=1)

        # Inpaint
        mask_pil = Image.fromarray(mask)
        with torch.no_grad():
            cleaned = self.lama(pil_image, mask_pil)

        # Cleanup
        del img_np, mask, mask_pil
        if self.is_gpu:
            torch.cuda.empty_cache()
        gc.collect()

        return cleaned, text_blocks


def _add_textbox_for_block(
    slide,
    block: dict,
    left_pt: float,
    top_pt: float,
    scale_x: float,
    scale_y: float,
) -> None:
    """Place one editable textbox on a slide for a single OCR block.

    Shared between the initial conversion path and the rebuild-from-
    workspace path so both produce identical box geometry.
    """
    box_left = Pt(left_pt + block["x_pt"] * scale_x)
    box_top = Pt(top_pt + block["y_pt"] * scale_y)
    raw_w = block["w_pt"] * scale_x
    raw_h = block["h_pt"] * scale_y
    box_w = Pt(max(raw_w * 1.12, 12.0))
    box_h = Pt(max(raw_h * 1.05, 10.0))
    txBox = slide.shapes.add_textbox(box_left, box_top, box_w, box_h)
    tf = txBox.text_frame
    # word_wrap=False keeps single-line text intact even when the OCR bbox
    # is slightly narrower than the rendered glyphs (so "LLM" doesn't wrap
    # to L/L/M). auto_size lets the shape grow horizontally instead.
    tf.word_wrap = False
    try:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass
    tf.text = block.get("text", "")
    if block.get("h_pt", 0) > 0:
        font_pt = max(6.0, raw_h * 0.65)
        tf.paragraphs[0].font.size = Pt(font_pt)


def _save_page_state(
    workspace_dir: Path,
    page_idx: int,
    cleaned_img: Image.Image,
    text_blocks: list[dict],
    page_w_pt: float,
    page_h_pt: float,
) -> None:
    """Persist (inpainted bg, blocks, page meta) so the user can review/edit
    later without re-running OCR + LaMa."""
    workspace_dir.mkdir(parents=True, exist_ok=True)
    bg_path = workspace_dir / f"page_{page_idx}_bg.png"
    cleaned_img.save(str(bg_path), format="PNG", optimize=True)
    json_path = workspace_dir / f"page_{page_idx}.json"
    json_path.write_text(
        json.dumps(
            {
                "page_idx": page_idx,
                "page_w_pt": page_w_pt,
                "page_h_pt": page_h_pt,
                "blocks": text_blocks,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )


def update_page_blocks(workspace_dir: Path, page_idx: int, blocks: list[dict]) -> bool:
    """Replace the saved blocks list for a page (after user edits)."""
    json_path = workspace_dir / f"page_{page_idx}.json"
    if not json_path.exists():
        return False
    data = json.loads(json_path.read_text(encoding="utf-8"))
    data["blocks"] = blocks
    json_path.write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return True


def load_page_state(workspace_dir: Path, page_idx: int) -> dict | None:
    json_path = workspace_dir / f"page_{page_idx}.json"
    if not json_path.exists():
        return None
    return json.loads(json_path.read_text(encoding="utf-8"))


def ocr_region_from_bg(
    workspace_dir: Path,
    page_idx: int,
    bbox_norm: list[float],
    ocr_engine,
) -> dict | None:
    """OCR a single user-drawn region of the (inpainted) background.

    The inpainted bg still contains text in regions OCR originally missed,
    so OCR'ing the user-marked crop is the right move.

    bbox_norm is [x, y, w, h] in 0..1 of the bg image. Returns a block dict
    in the same shape as the auto-detected blocks, or None if nothing is
    recognised.
    """
    json_path = workspace_dir / f"page_{page_idx}.json"
    bg_path = workspace_dir / f"page_{page_idx}_bg.png"
    if not json_path.exists() or not bg_path.exists():
        return None
    data = json.loads(json_path.read_text(encoding="utf-8"))
    page_w_pt = float(data["page_w_pt"])

    img = Image.open(bg_path).convert("RGB")
    W, H = img.size
    x_norm, y_norm, w_norm, h_norm = bbox_norm
    left_px = max(0, int(x_norm * W))
    top_px = max(0, int(y_norm * H))
    right_px = min(W, int((x_norm + w_norm) * W))
    bottom_px = min(H, int((y_norm + h_norm) * H))
    if right_px <= left_px or bottom_px <= top_px:
        return None

    crop = img.crop((left_px, top_px, right_px, bottom_px))
    crop_np = np.array(crop)
    try:
        ocr_result, _ = ocr_engine(crop_np)
    except Exception:
        ocr_result = None

    texts = []
    if ocr_result:
        for item in ocr_result:
            try:
                _box, t, _s = item
            except (ValueError, TypeError):
                continue
            if t and t.strip():
                texts.append(t.strip())
    combined = " ".join(texts).strip()
    if not combined:
        return None

    px_per_pt = max(W / page_w_pt, 1e-6)
    return {
        "text": combined,
        "x_pt": left_px / px_per_pt,
        "y_pt": top_px / px_per_pt,
        "w_pt": (right_px - left_px) / px_per_pt,
        "h_pt": (bottom_px - top_px) / px_per_pt,
        "score": 0.9,
    }


def assemble_pptx_from_workspace(workspace_dir: Path, output_path: Path) -> Path:
    """Rebuild a PPTX from the per-page state saved during the original
    conversion (plus any edits the user has made via the review API)."""
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = prs.slide_layouts[6]

    json_files = sorted(
        workspace_dir.glob("page_*.json"),
        key=lambda p: int(p.stem.split("_")[1]),
    )
    for jf in json_files:
        try:
            data = json.loads(jf.read_text(encoding="utf-8"))
        except Exception:
            continue
        page_idx = int(data.get("page_idx", -1))
        if page_idx < 0:
            continue
        bg_path = workspace_dir / f"page_{page_idx}_bg.png"
        if not bg_path.exists():
            continue
        page_w_pt = float(data["page_w_pt"])
        page_h_pt = float(data["page_h_pt"])
        blocks = data.get("blocks", []) or []

        slide = prs.slides.add_slide(blank_layout)
        left_pt, top_pt, fit_w_pt, fit_h_pt = _fit_image_to_slide(page_w_pt, page_h_pt)
        scale_x = fit_w_pt / page_w_pt
        scale_y = fit_h_pt / page_h_pt

        slide.shapes.add_picture(
            str(bg_path),
            Pt(left_pt),
            Pt(top_pt),
            width=Pt(fit_w_pt),
            height=Pt(fit_h_pt),
        )
        for block in blocks:
            try:
                _add_textbox_for_block(slide, block, left_pt, top_pt, scale_x, scale_y)
            except Exception:
                continue

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _fit_image_to_slide(
    page_w_pt: float, page_h_pt: float
) -> tuple[float, float, float, float]:
    """Fix #1: fit page image into 16:9 slide (letterbox). Returns (left_pt, top_pt, w_pt, h_pt)."""
    page_aspect = page_w_pt / page_h_pt
    slide_aspect = SLIDE_W_PT / SLIDE_H_PT
    if abs(page_aspect - slide_aspect) < 0.01:
        return 0.0, 0.0, SLIDE_W_PT, SLIDE_H_PT
    if page_aspect > slide_aspect:
        # Page is wider -> letterbox top/bottom
        new_w = SLIDE_W_PT
        new_h = SLIDE_W_PT / page_aspect
        return 0.0, (SLIDE_H_PT - new_h) / 2, new_w, new_h
    # Page is taller -> letterbox left/right
    new_h = SLIDE_H_PT
    new_w = SLIDE_H_PT * page_aspect
    return (SLIDE_W_PT - new_w) / 2, 0.0, new_w, new_h


def convert_pdf_to_pptx(
    pdf_path: Path,
    output_path: Path,
    dpi: int = 200,
    dilation: int = 15,
    converter: Optional[LocalConverter] = None,
    on_page_progress: Optional[Callable[[int, str, dict], None]] = None,
    cancel_flag: Optional[Callable[[], bool]] = None,
    workspace_dir: Optional[Path] = None,
) -> Path:
    """Convert one PDF into an editable PPTX.

    on_page_progress(page_idx, state, info) is called as the pipeline moves;
    state is one of: "rendering", "ocr", "inpainting", "done", "error".
    info carries extra detail (text_block_count, elapsed_s, message).
    """
    if converter is None:
        converter = LocalConverter()
        converter.initialize_models()

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = prs.slide_layouts[6]

    # pypdfium2 (pdfplumber's backend) is not thread-safe. Hold the lock
    # for the whole conversion so any concurrent thumbnail/conversion call
    # waits its turn instead of crashing in native code.
    with _pdf_lock, pdfplumber.open(str(pdf_path)) as pdf_file:
        for i, page in enumerate(pdf_file.pages):
            if cancel_flag and cancel_flag():
                raise RuntimeError("Cancelled by user")

            t0 = time.monotonic()
            page_w_pt = float(page.width)
            page_h_pt = float(page.height)

            if on_page_progress:
                on_page_progress(i, "rendering", {})

            try:
                page_image_obj = page.to_image(resolution=dpi)
                pil_image = page_image_obj.original.convert("RGB")
            except Exception as exc:
                if on_page_progress:
                    on_page_progress(i, "error", {"message": f"render failed: {exc}"})
                continue

            if on_page_progress:
                on_page_progress(i, "ocr", {})
            try:
                cleaned, text_blocks = converter.process_page(pil_image, dpi, dilation)
            except Exception as exc:
                if on_page_progress:
                    on_page_progress(i, "error", {"message": f"OCR/inpaint failed: {exc}"})
                continue

            if on_page_progress:
                on_page_progress(i, "inpainting", {"text_block_count": len(text_blocks)})

            # Persist per-page state for review/regenerate, before assembling
            # the slide. Doing it here means even if PPTX assembly fails later
            # we still have the pieces.
            if workspace_dir is not None:
                try:
                    _save_page_state(
                        workspace_dir, i, cleaned, text_blocks, page_w_pt, page_h_pt
                    )
                except Exception:
                    pass

            slide = prs.slides.add_slide(blank_layout)

            # Fix #1: letterbox the page image into the 16:9 slide
            left_pt, top_pt, fit_w_pt, fit_h_pt = _fit_image_to_slide(
                page_w_pt, page_h_pt
            )
            scale_x = fit_w_pt / page_w_pt
            scale_y = fit_h_pt / page_h_pt

            bg_buf = BytesIO()
            cleaned.save(bg_buf, format="PNG", optimize=True)
            bg_buf.seek(0)
            slide.shapes.add_picture(
                bg_buf,
                Pt(left_pt),
                Pt(top_pt),
                width=Pt(fit_w_pt),
                height=Pt(fit_h_pt),
            )

            for block in text_blocks:
                try:
                    _add_textbox_for_block(slide, block, left_pt, top_pt, scale_x, scale_y)
                except Exception:
                    continue

            elapsed = time.monotonic() - t0
            if on_page_progress:
                on_page_progress(
                    i,
                    "done",
                    {"text_block_count": len(text_blocks), "elapsed_s": round(elapsed, 1)},
                )

            del pil_image, cleaned, page_image_obj
            gc.collect()

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# Module-level shared converter (models loaded lazily on first request)
_shared_converter: Optional[LocalConverter] = None
_shared_lock = threading.Lock()


def get_shared_converter() -> LocalConverter:
    global _shared_converter
    with _shared_lock:
        if _shared_converter is None:
            _shared_converter = LocalConverter()
        return _shared_converter
