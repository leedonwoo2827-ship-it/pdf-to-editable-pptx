"""Compose the .pptx output: 16:9 slides with letterboxed page bg + text boxes.

Two callers:
- The conversion pipeline writes slides as it processes each page.
- The review API rebuilds a fresh .pptx from saved workspace state.

Both share the same slide geometry so a re-export looks identical to
the original conversion.
"""
from __future__ import annotations

from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt


# 16:9 widescreen at PowerPoint default scale. EMU = 914400 per inch.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SLIDE_W_EMU = Emu(int(SLIDE_W_IN * 914400))
SLIDE_H_EMU = Emu(int(SLIDE_H_IN * 914400))
SLIDE_W_PT = SLIDE_W_IN * 72
SLIDE_H_PT = SLIDE_H_IN * 72


def new_presentation() -> Presentation:
    """Build an empty 16:9 Presentation with our standard size."""
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    return prs


def fit_image_to_slide(
    page_w_pt: float, page_h_pt: float
) -> tuple[float, float, float, float]:
    """Letterbox a page image into the 16:9 slide. Returns
    (left_pt, top_pt, width_pt, height_pt) for the picture placement."""
    page_aspect = page_w_pt / page_h_pt
    slide_aspect = SLIDE_W_PT / SLIDE_H_PT
    if abs(page_aspect - slide_aspect) < 0.01:
        return 0.0, 0.0, SLIDE_W_PT, SLIDE_H_PT
    if page_aspect > slide_aspect:
        # Page is wider than slide → letterbox top/bottom.
        new_w = SLIDE_W_PT
        new_h = SLIDE_W_PT / page_aspect
        return 0.0, (SLIDE_H_PT - new_h) / 2, new_w, new_h
    # Page is taller than slide → letterbox left/right.
    new_h = SLIDE_H_PT
    new_w = SLIDE_H_PT * page_aspect
    return (SLIDE_W_PT - new_w) / 2, 0.0, new_w, new_h


def add_textbox_for_block(
    slide,
    block: dict,
    left_pt: float,
    top_pt: float,
    scale_x: float,
    scale_y: float,
) -> None:
    """Place one editable textbox on the slide for a single block.

    word_wrap=False keeps single-line tight text from breaking glyph-
    by-glyph when the OCR bbox is slightly narrower than the rendered
    text. SHAPE_TO_FIT_TEXT lets the shape grow horizontally instead.
    """
    box_left = Pt(left_pt + block["x_pt"] * scale_x)
    box_top = Pt(top_pt + block["y_pt"] * scale_y)
    raw_w = block["w_pt"] * scale_x
    raw_h = block["h_pt"] * scale_y
    box_w = Pt(max(raw_w * 1.12, 12.0))
    box_h = Pt(max(raw_h * 1.05, 10.0))

    txBox = slide.shapes.add_textbox(box_left, box_top, box_w, box_h)
    tf = txBox.text_frame
    tf.word_wrap = False
    try:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass
    tf.text = block.get("text", "")
    if block.get("h_pt", 0) > 0:
        # Empirical: bbox_h * 0.65 gives a good visual font size match
        # at typical DPI. Floor at 6pt so degenerate tiny boxes still
        # render readable text.
        font_pt = max(6.0, raw_h * 0.65)
        tf.paragraphs[0].font.size = Pt(font_pt)


def add_page_slide(
    prs: Presentation,
    cleaned_bg: Image.Image | bytes | str,
    text_blocks: list[dict],
    page_w_pt: float,
    page_h_pt: float,
) -> None:
    """Add one slide with the given background image + text blocks.

    cleaned_bg can be a PIL.Image (will be PNG-encoded into memory), a
    bytes object (assumed PNG), or a str/Path-like (read from disk).
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    left_pt, top_pt, fit_w_pt, fit_h_pt = fit_image_to_slide(page_w_pt, page_h_pt)
    scale_x = fit_w_pt / page_w_pt
    scale_y = fit_h_pt / page_h_pt

    if isinstance(cleaned_bg, Image.Image):
        buf = BytesIO()
        cleaned_bg.save(buf, format="PNG", optimize=True)
        buf.seek(0)
        picture_arg: object = buf
    else:
        picture_arg = str(cleaned_bg) if not isinstance(cleaned_bg, (bytes, bytearray)) else BytesIO(cleaned_bg)

    slide.shapes.add_picture(
        picture_arg,
        Pt(left_pt),
        Pt(top_pt),
        width=Pt(fit_w_pt),
        height=Pt(fit_h_pt),
    )

    for block in text_blocks:
        try:
            add_textbox_for_block(slide, block, left_pt, top_pt, scale_x, scale_y)
        except Exception:
            # One bad block shouldn't kill the whole slide.
            continue
